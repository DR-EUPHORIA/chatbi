"""文件导出工具 — 生成 HTML 报告、Excel、Markdown、PPT"""

import json
from pathlib import Path
from typing import Any

import pandas as pd
from jinja2 import Template


# ── HTML 报告模板 ──
HTML_REPORT_TEMPLATE = """<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{ title }}</title>
    <script src="https://cdn.jsdelivr.net/npm/echarts@5/dist/echarts.min.js"></script>
    <style>
        * { margin: 0; padding: 0; box-sizing: border-box; }
        body {
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, 'PingFang SC', 'Microsoft YaHei', sans-serif;
            background: #f5f7fa;
            color: #333;
            line-height: 1.8;
        }
        .container { max-width: 1000px; margin: 0 auto; padding: 40px 20px; }
        .header {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 40px;
            border-radius: 16px;
            margin-bottom: 30px;
        }
        .header h1 { font-size: 28px; margin-bottom: 12px; }
        .header p { font-size: 16px; opacity: 0.9; }
        .card {
            background: white;
            border-radius: 12px;
            padding: 30px;
            margin-bottom: 20px;
            box-shadow: 0 2px 12px rgba(0,0,0,0.08);
        }
        .card h2 {
            font-size: 20px;
            color: #667eea;
            margin-bottom: 16px;
            padding-bottom: 8px;
            border-bottom: 2px solid #f0f0f0;
        }
        .chart-container { width: 100%; height: 400px; }
        .insight-list { list-style: none; }
        .insight-list li {
            padding: 12px 16px;
            margin-bottom: 8px;
            background: #f8f9ff;
            border-left: 4px solid #667eea;
            border-radius: 0 8px 8px 0;
        }
        .recommendation-list li {
            border-left-color: #52c41a;
            background: #f6ffed;
        }
        .footer {
            text-align: center;
            padding: 20px;
            color: #999;
            font-size: 14px;
        }
        .generated-by { color: #667eea; font-weight: 600; }
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>{{ title }}</h1>
            <p>{{ summary }}</p>
        </div>

        {% for section in sections %}
        <div class="card">
            <h2>{{ section.section_title }}</h2>
            <div>{{ section.content }}</div>
        </div>
        {% endfor %}

        {% if chart_config %}
        <div class="card">
            <h2>数据可视化</h2>
            <div id="main-chart" class="chart-container"></div>
        </div>
        {% endif %}

        {% if insights %}
        <div class="card">
            <h2>关键洞察</h2>
            <ul class="insight-list">
                {% for insight in insights %}
                <li>{{ insight }}</li>
                {% endfor %}
            </ul>
        </div>
        {% endif %}

        {% if recommendations %}
        <div class="card">
            <h2>行动建议</h2>
            <ul class="insight-list recommendation-list">
                {% for rec in recommendations %}
                <li>{{ rec }}</li>
                {% endfor %}
            </ul>
        </div>
        {% endif %}

        <div class="footer">
            由 <span class="generated-by">ChatBI Mini</span> 自动生成
        </div>
    </div>

    {% if chart_config %}
    <script>
        var chart = echarts.init(document.getElementById('main-chart'));
        var option = {{ chart_config_json }};
        chart.setOption(option);
        window.addEventListener('resize', function() { chart.resize(); });
    </script>
    {% endif %}
</body>
</html>"""


def generate_html_report(
    title: str,
    summary: str,
    sections: list[dict],
    insights: list[str],
    recommendations: list[str],
    chart_config: dict,
    output_path: str,
) -> str:
    """生成 HTML 报告文件"""
    template = Template(HTML_REPORT_TEMPLATE)
    html_content = template.render(
        title=title,
        summary=summary,
        sections=sections,
        insights=insights,
        recommendations=recommendations,
        chart_config=chart_config,
        chart_config_json=json.dumps(chart_config, ensure_ascii=False),
    )

    Path(output_path).write_text(html_content, encoding="utf-8")
    return output_path


def generate_excel_export(
    data: list[dict],
    columns: list[str],
    analysis: dict[str, Any],
    output_path: str,
) -> str:
    """生成 Excel 导出文件（多 Sheet）"""
    with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
        # Sheet 1: 原始数据
        if data:
            dataframe = pd.DataFrame(data, columns=columns)
            dataframe.to_excel(writer, sheet_name="数据明细", index=False)

        # Sheet 2: 分析摘要
        summary_data = []
        key_metrics = analysis.get("key_metrics", {})
        for metric_name, metric_value in key_metrics.items():
            summary_data.append({"指标": metric_name, "值": str(metric_value)})

        if summary_data:
            summary_df = pd.DataFrame(summary_data)
            summary_df.to_excel(writer, sheet_name="关键指标", index=False)

        # Sheet 3: 洞察
        insights = analysis.get("insights", [])
        if insights:
            insights_df = pd.DataFrame({"序号": range(1, len(insights) + 1), "洞察": insights})
            insights_df.to_excel(writer, sheet_name="数据洞察", index=False)

    return output_path


def generate_ppt_export(
    title: str,
    summary: str,
    insights: list[str],
    recommendations: list[str],
    output_path: str,
) -> str:
    """生成 PPT 演示文稿"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    presentation = Presentation()
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    # ── 封面页 ──
    slide_layout = presentation.slide_layouts[6]  # 空白布局
    slide = presentation.slides.add_slide(slide_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0x66, 0x7E, 0xEA)
    title_para.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = summary
    subtitle_para.font.size = Pt(16)
    subtitle_para.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 来源标注
    source_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    source_frame = source_box.text_frame
    source_para = source_frame.paragraphs[0]
    source_para.text = "由 ChatBI Mini 自动生成"
    source_para.font.size = Pt(12)
    source_para.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    source_para.alignment = PP_ALIGN.CENTER

    # ── 洞察页 ──
    if insights:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])

        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        header_frame = header_box.text_frame
        header_para = header_frame.paragraphs[0]
        header_para.text = "关键洞察"
        header_para.font.size = Pt(28)
        header_para.font.bold = True
        header_para.font.color.rgb = RGBColor(0x66, 0x7E, 0xEA)

        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for idx, insight in enumerate(insights):
            para = content_frame.add_paragraph() if idx > 0 else content_frame.paragraphs[0]
            para.text = f"  {idx + 1}. {insight}"
            para.font.size = Pt(14)
            para.space_after = Pt(12)
            para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ── 建议页 ──
    if recommendations:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])

        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        header_frame = header_box.text_frame
        header_para = header_frame.paragraphs[0]
        header_para.text = "行动建议"
        header_para.font.size = Pt(28)
        header_para.font.bold = True
        header_para.font.color.rgb = RGBColor(0x52, 0xC4, 0x1A)

        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for idx, rec in enumerate(recommendations):
            para = content_frame.add_paragraph() if idx > 0 else content_frame.paragraphs[0]
            para.text = f"  {idx + 1}. {rec}"
            para.font.size = Pt(14)
            para.space_after = Pt(12)
            para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    presentation.save(output_path)
    return output_path


def generate_export_file(session_id: str, export_format: str, content: dict) -> dict:
    """通用导出接口"""
    from config import settings

    export_dir = Path(settings.export_dir)
    export_dir.mkdir(parents=True, exist_ok=True)

    if export_format == "html":
        output_path = str(export_dir / f"report_{session_id}.html")
        generate_html_report(
            title=content.get("title", "数据分析报告"),
            summary=content.get("summary", ""),
            sections=content.get("sections", []),
            insights=content.get("insights", []),
            recommendations=content.get("recommendations", []),
            chart_config=content.get("chart_config", {}),
            output_path=output_path,
        )
        return {"url": f"/exports/report_{session_id}.html", "format": "html"}

    elif export_format == "excel":
        output_path = str(export_dir / f"data_{session_id}.xlsx")
        generate_excel_export(
            data=content.get("data", []),
            columns=content.get("columns", []),
            analysis=content.get("analysis", {}),
            output_path=output_path,
        )
        return {"url": f"/exports/data_{session_id}.xlsx", "format": "excel"}

    elif export_format == "markdown":
        output_path = str(export_dir / f"report_{session_id}.md")
        Path(output_path).write_text(content.get("markdown", ""), encoding="utf-8")
        return {"url": f"/exports/report_{session_id}.md", "format": "markdown"}

    elif export_format == "ppt":
        output_path = str(export_dir / f"report_{session_id}.pptx")
        generate_ppt_export(
            title=content.get("title", "数据分析报告"),
            summary=content.get("summary", ""),
            insights=content.get("insights", []),
            recommendations=content.get("recommendations", []),
            output_path=output_path,
        )
        return {"url": f"/exports/report_{session_id}.pptx", "format": "ppt"}

    else:
        return {"error": f"不支持的导出格式：{export_format}"}
